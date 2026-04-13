import sys, io
# Force UTF-8 on Windows to avoid cp1252 UnicodeEncodeError
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
else:
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

import os
import re
import math
import platform
import subprocess
import psutil
from pathlib import Path
from datetime import datetime
from collections import defaultdict

# ── pyfiglet (ASCII art titles) ───────────────────────────────────────────────
try:
    import pyfiglet
    FIGLET_OK = True
except ImportError:
    FIGLET_OK = False

# ── output library imports (graceful) ─────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    from docx import Document
    from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    from llama_cpp import Llama
    LLAMA_OK = True
except ImportError:
    LLAMA_OK = False

# ── win32com (COM automation – Windows only) ───────────────────────────────────
try:
    import win32com.client
    import pythoncom
    WIN32_OK = True
except ImportError:
    WIN32_OK = False

# ── paths ──────────────────────────────────────────────────────────────────────
ROOT       = Path(__file__).parent
WORKSPACE  = ROOT / "workspace"
OUTPUTS    = ROOT / "outputs"
WORKSPACE.mkdir(exist_ok=True)
OUTPUTS.mkdir(exist_ok=True)

# ── supported extensions ───────────────────────────────────────────────────────
TEXT_EXTS  = {".txt", ".md", ".rst", ".csv", ".log"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp"}
# NEW: workspace document types that can be read and summarised
READABLE_DOC_EXTS = {".pptx", ".docx", ".txt", ".md", ".rst", ".csv", ".log"}

# ── layout constants ───────────────────────────────────────────────────────────
W     = 68
LEFT  = 4
PAD   = " " * LEFT
INNER = W - LEFT - 4   # usable text width inside box walls  (= 60)

# ── PaperReady identity ────────────────────────────────────────────────────────
IDENTITY = {
    "name"       : "PaperReady",
    "designer"   : "Bhargav",
    "base_model" : "Phi-3 Mini 4K Instruct by Microsoft",
    "version"    : "2.0",
    "purpose"    : (
        "PaperReady is an AI-powered document and presentation assistant "
        "designed and built by Bhargav. It runs a local Phi-3 Mini language model "
        "enhanced with workspace-aware RAG (Retrieval-Augmented Generation). "
        "PaperReady can create rich PowerPoint presentations, formatted Word documents, "
        "and structured text notes — written to disk and opened directly in PowerPoint "
        "or Word on your PC. It can also embed your own images from the workspace/ folder."
    ),
    "capabilities": [
        "Chat with a local Phi-3 AI — no internet needed",
        "Create PowerPoint presentations via live COM automation (slides written inside PPT)",
        "Create formatted Word (.docx) documents via COM automation",
        "Write plain-text notes (.txt) with auto-header",
        "Insert workspace images into slides and documents automatically",
        "RAG over your workspace/ folder for topic-aware generation",
        "Read and summarise .pptx, .docx, and .txt files from workspace/",
        "Reload workspace index at any time with 'reload'",
        "Tell you about itself when asked",
    ],
}

ABOUT_TRIGGERS = re.compile(
    r"\b(who are you|what are you|about you|tell me about yourself|"
    r"what can you do|your capabilities|what is paperready|"
    r"who (made|built|designed|created) you|your (name|creator|designer|model))\b",
    re.I,
)

# NEW: pattern to detect "read/summarise file X" requests
SUMMARISE_TRIGGERS = re.compile(
    r"\b(summarize|summarise|summary of|read|open|explain|describe|"
    r"what('?s| is) in|contents? of|give me a summary)\b.{0,60}"
    r"(\.pptx|\.docx|\.txt|\.md|\.csv|\.log|\.rst)\b",
    re.I,
)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 0 — ASCII Art Helpers
# ══════════════════════════════════════════════════════════════════════════════

def _hr(char: str = "=") -> str:
    return PAD + char * (W - LEFT)

def _box_line(text: str = "", char: str = "|") -> str:
    """Single box line. Text is ALWAYS clipped to INNER width so the right wall stays straight."""
    if len(text) > INNER:
        text = text[:INNER - 3] + "..."
    return PAD + char + " " + text.ljust(INNER) + " " + char

def _box_wrap_lines(text: str, indent: str = "", char: str = "|") -> list:
    """Word-wrap text so every printed line fits inside the box walls."""
    available = INNER - len(indent)
    words     = text.split()
    rows      = []
    line = ""
    for word in words:
        candidate = f"{line} {word}".strip() if line else word
        if len(candidate) > available:
            if line:
                rows.append(_box_line(indent + line, char))
            line = word
        else:
            line = candidate
    if line:
        rows.append(_box_line(indent + line, char))
    return rows if rows else [_box_line(indent, char)]

def _center_line(text: str, char: str = "|") -> str:
    return PAD + char + " " + text.center(INNER) + " " + char

def _label(tag: str, text: str) -> str:
    tag_str = f"[{tag:^7}]"
    return f"{PAD}{tag_str}  {text}"

def print_section(title: str):
    print()
    print(_hr("-"))
    print(_label("  >>  ", title))
    print(_hr("-"))

def print_status(symbol: str, message: str):
    print(f"{PAD}  {symbol}  {message}")

def print_info(message: str):
    print(f"{PAD}  *  {message}")

def print_ok(message: str):
    print(f"{PAD}  [OK]  {message}")

def print_err(message: str):
    print(f"{PAD}  [!!]  {message}")

def _figlet_banner(text: str, font: str = "slant") -> str:
    if FIGLET_OK:
        try:
            raw   = pyfiglet.figlet_format(text, font=font)
            lines = raw.splitlines()
            return "\n".join(PAD + line for line in lines)
        except Exception:
            pass
    border = PAD + "+" + "-" * (W - LEFT - 2) + "+"
    inner  = W - LEFT - 4
    middle = PAD + "|" + text.center(inner + 2) + "|"
    return "\n".join([border, middle, border])


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — Platform-Aware File Opener
# ══════════════════════════════════════════════════════════════════════════════

def ask_permission(action: str) -> bool:
    print()
    print(f"{PAD}  [PERMISSION]  {action}")
    print(f"{PAD}  Allow? (yes / no)  ->  ", end="", flush=True)
    try:
        ans = input().strip().lower()
    except (EOFError, KeyboardInterrupt):
        print()
        return False
    return ans in ("yes", "y", "1", "ok", "sure", "yep", "yeah")


def open_file_in_app(file_path: Path, app_hint: str = "auto") -> bool:
    system   = platform.system()
    path_str = str(file_path.resolve())
    try:
        if system == "Windows":
            if app_hint == "txt":
                subprocess.Popen(["notepad.exe", path_str])
            else:
                os.startfile(path_str)
        elif system == "Darwin":
            subprocess.Popen(["open", path_str])
        else:
            subprocess.Popen(["xdg-open", path_str])
        return True
    except Exception as e:
        print_err(f"Could not open file automatically: {e}")
        print_info(f"Please open manually:  {path_str}")
        return False


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 2 — Workspace RAG
# ══════════════════════════════════════════════════════════════════════════════

class WorkspaceRAG:
    def __init__(self):
        self.text_docs: dict  = {}
        self.image_paths: list = []
        self.tfidf: dict = {}
        self._build_index()

    def _build_index(self):
        for p in WORKSPACE.rglob("*"):
            if p.is_file():
                if p.suffix.lower() in TEXT_EXTS:
                    try:
                        self.text_docs[p.name] = p.read_text(encoding="utf-8", errors="ignore")
                    except Exception:
                        pass
                elif p.suffix.lower() in IMAGE_EXTS:
                    self.image_paths.append(p)
                # NEW: index .pptx and .docx as extracted text too
                elif p.suffix.lower() == ".pptx":
                    extracted = _extract_text_from_pptx(p)
                    if extracted:
                        self.text_docs[p.name] = extracted
                elif p.suffix.lower() == ".docx":
                    extracted = _extract_text_from_docx(p)
                    if extracted:
                        self.text_docs[p.name] = extracted
        if self.text_docs:
            self._compute_tfidf()
        print_info(
            f"RAG index built  ->  {len(self.text_docs)} text/doc file(s), "
            f"{len(self.image_paths)} image(s) in workspace/"
        )

    def _tokenize(self, text: str) -> list:
        return re.findall(r"[a-z]+", text.lower())

    def _compute_tfidf(self):
        tf  = defaultdict(lambda: defaultdict(int))
        df  = defaultdict(int)
        N = len(self.text_docs)
        for name, content in self.text_docs.items():
            words = self._tokenize(content)
            for w in words:
                tf[name][w] += 1
            for w in set(words):
                df[w] += 1
        for name, word_counts in tf.items():
            total = sum(word_counts.values()) or 1
            for w, count in word_counts.items():
                idf = math.log((N + 1) / (df[w] + 1)) + 1
                if w not in self.tfidf:
                    self.tfidf[w] = {}
                self.tfidf[w][name] = (count / total) * idf

    def retrieve_text(self, query: str, top_k: int = 3, max_chars: int = 1200) -> str:
        if not self.text_docs:
            return ""
        q_words = set(self._tokenize(query))
        scores  = defaultdict(float)
        for w in q_words:
            if w in self.tfidf:
                for docname, score in self.tfidf[w].items():
                    scores[docname] += score
        ranked   = sorted(scores.items(), key=lambda x: x[1], reverse=True)
        snippets = []
        for docname, _ in ranked[:top_k]:
            text = self.text_docs[docname]
            snippets.append(f"[From: {docname}]\n{text[:max_chars]}")
        return "\n\n".join(snippets)

    def retrieve_images(self, query: str) -> list:
        q_words  = set(self._tokenize(query))
        q_phrase = re.sub(r"\s+", "_", query.strip().lower())
        matched  = []
        for img in self.image_paths:
            stem_words = set(self._tokenize(img.stem))
            if q_words & stem_words or q_phrase in img.stem.lower():
                matched.append(img)
        return matched

    def reload(self):
        self.text_docs.clear()
        self.image_paths.clear()
        self.tfidf.clear()
        self._build_index()

    # NEW: find a workspace file by filename (partial match)
    def find_file(self, filename_hint: str) -> Path | None:
        hint_lower = filename_hint.strip().lower()
        for p in WORKSPACE.rglob("*"):
            if p.is_file() and p.suffix.lower() in READABLE_DOC_EXTS:
                if hint_lower in p.name.lower():
                    return p
        return None


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 2b — Workspace File Text Extractors  (NEW)
# ══════════════════════════════════════════════════════════════════════════════

def _extract_text_from_pptx(path: Path) -> str:
    """Extract all text from a .pptx file using python-pptx."""
    if not PPTX_OK:
        return ""
    try:
        prs    = Presentation(str(path))
        lines  = [f"[PowerPoint: {path.name}]"]
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt:
                            lines.append(txt)
        return "\n".join(lines)
    except Exception as e:
        print_err(f"Could not read {path.name}: {e}")
        return ""


def _extract_text_from_docx(path: Path) -> str:
    """Extract all text from a .docx file using python-docx."""
    if not DOCX_OK:
        return ""
    try:
        doc   = Document(str(path))
        lines = [f"[Word Document: {path.name}]"]
        for para in doc.paragraphs:
            txt = para.text.strip()
            if txt:
                lines.append(txt)
        return "\n".join(lines)
    except Exception as e:
        print_err(f"Could not read {path.name}: {e}")
        return ""


def _extract_text_from_txt(path: Path) -> str:
    """Read a plain text file."""
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        print_err(f"Could not read {path.name}: {e}")
        return ""


def extract_workspace_file(path: Path) -> str:
    """Route to the correct extractor based on file extension."""
    ext = path.suffix.lower()
    if ext == ".pptx":
        return _extract_text_from_pptx(path)
    elif ext == ".docx":
        return _extract_text_from_docx(path)
    else:
        return _extract_text_from_txt(path)


def detect_summarise_request(query: str) -> str | None:
    """
    If the query looks like 'summarize myfile.pptx', extract and return the
    filename hint (including extension).  Returns None if no match.
    """
    # Match patterns like: "summarize report.pptx", "read notes.txt", "what's in doc.docx"
    m = re.search(
        r"([\w\-\.]+\.(pptx|docx|txt|md|csv|log|rst))",
        query,
        re.I,
    )
    if m:
        return m.group(1)
    return None


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 3 — Intent Detection & Slide Count Extraction
# ══════════════════════════════════════════════════════════════════════════════

_INTENT_PATTERNS = {
    "pptx": re.compile(r"\b(powerpoint|pptx|presentation|slides?|slide deck)\b", re.I),
    "docx": re.compile(r"\b(word\s+doc(ument)?|docx|\.docx)\b", re.I),
    "txt" : re.compile(r"\b(notepad|text\s*file|\.txt|write\s+(to\s+)?file|save\s+(as\s+)?text)\b", re.I),
    # NEW: summarise intent
    "summarise": re.compile(
        r"\b(summarize|summarise|summary|read|open|explain|describe|"
        r"what('?s| is) in|contents? of|give me a summary)\b",
        re.I,
    ),
}

def detect_intent(query: str) -> str:
    # Check for summarise FIRST (before generic chat), only if a filename is present
    if _INTENT_PATTERNS["summarise"].search(query) and detect_summarise_request(query):
        return "summarise"
    for intent, pattern in _INTENT_PATTERNS.items():
        if intent == "summarise":
            continue   # already handled above
        if pattern.search(query):
            return intent
    return "chat"

def extract_topic(query: str) -> str:
    q = query.strip()
    q = re.sub(r"\b\d+[\s-]*slides?\b", " ", q, flags=re.I)
    for kw in ["powerpoint", "pptx", "ppt", "presentation", "slides", "slide deck",
               "word document", "word doc", "docx", "notepad", "text file",
               "write a", "create a", "make a", "generate a",
               "write", "create", "make", "generate",
               "about", "on", "for", "regarding", "please", "me", "a"]:
        q = re.sub(r"\b" + re.escape(kw) + r"\b", " ", q, flags=re.I)
    topic = re.sub(r"\s+", " ", q).strip().strip(".,;:'\"")
    topic = re.sub(
        r"\s+\b(in|on|for|the|an|a|of|and|or|with|by|to|into|from|about|at|as)\b\s*$",
        "", topic, flags=re.I
    ).strip().strip(".,;:'\"")
    return topic if topic else "My Topic"

def extract_slide_count(query: str) -> int:
    word_map = {
        "one":1,"two":2,"three":3,"four":4,"five":5,
        "six":6,"seven":7,"eight":8,"nine":9,"ten":10,
    }
    m = re.search(r"\b(\d+)[\s-]*slides?\b", query, re.I)
    if m:
        return max(2, min(20, int(m.group(1))))
    for word, num in word_map.items():
        if re.search(r"\b" + word + r"[\s-]*slides?\b", query, re.I):
            return num
    return 6


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 4 — LLM Prompt Builder
# ══════════════════════════════════════════════════════════════════════════════

def _identity_block() -> str:
    return (
        f"Your name is {IDENTITY['name']}. "
        f"You were designed and built by {IDENTITY['designer']}. "
        f"Your base language model is {IDENTITY['base_model']}. "
        "When asked who you are, introduce yourself accordingly and list your capabilities."
    )

SYSTEM_PROMPT_BASE = (
    "You are PaperReady, an expert writing and presentation assistant "
    "designed by Bhargav, powered by Phi-3 Mini by Microsoft. "
    "When given retrieved context, use it to improve accuracy.\n\n"

    "====== STRICT OUTPUT FORMAT FOR PRESENTATIONS ======\n"
    "CRITICAL RULE: When creating a presentation, output ONLY the SLIDE/CONTENT lines. "
    "Do NOT write any greeting, explanation, preamble, or text before SLIDE 1. "
    "Do NOT write anything after the last CONTENT line. "
    "Start your response with 'SLIDE 1:' — nothing before it, not even a blank line.\n\n"
    "Each slide: one SLIDE line followed immediately by one CONTENT line.\n"
    "CONTENT must have 3-5 points separated by ' | ' (pipe with spaces).\n\n"
    "EXAMPLE — 3-slide presentation about Solar Energy:\n"
    "SLIDE 1: Solar Energy — Powering the Future\n"
    "CONTENT: Clean renewable energy from the sun | 173,000 terawatts hit Earth daily | Zero greenhouse gas emissions | Fastest-growing energy source globally\n"
    "SLIDE 2: How Solar Panels Work\n"
    "CONTENT: Photovoltaic cells convert sunlight to electricity | Silicon semiconductors absorb photons | Direct current converted to AC by inverter | Excess power stored in batteries or fed to grid\n"
    "SLIDE 3: Benefits and the Road Ahead\n"
    "CONTENT: Reduces electricity bills by up to 70% | 25-30 year panel lifespan | Costs have fallen 90% since 2010 | Expected to supply 20% of global power by 2040\n\n"
    "IMPORTANT: Your FIRST character must be 'S' (start of SLIDE 1). No other text allowed before or after.\n\n"
    "====== FORMAT FOR WORD DOCUMENTS ======\n"
    "Use HEADING: markers for section titles, then prose paragraphs below each heading. "
    "Include at least 4 sections. No markdown code fences.\n\n"
    "====== FORMAT FOR TEXT NOTES ======\n"
    "Plain text, clear section headings, no formatting syntax.\n\n"
    + _identity_block()
)

def build_messages(history: list, query: str, context: str, intent: str,
                   topic: str, slide_count: int = 6) -> list:
    system_content = SYSTEM_PROMPT_BASE
    if context:
        system_content += f"\n\n--- Workspace Knowledge ---\n{context}"

    if intent == "pptx":
        task_prefix = (
            f"Create a PowerPoint presentation about '{topic}'. "
            f"Generate EXACTLY {slide_count} slides using the SLIDE/CONTENT format shown above. "
            f"Slide 1 = title/intro slide. "
            f"Slides 2 to {slide_count - 1} = key topic slides. "
            f"Slide {slide_count} = conclusion or summary. "
            "Output ONLY the SLIDE/CONTENT lines. Nothing else."
        )
    elif intent == "docx":
        task_prefix = (
            f"Write a professional Word document about '{topic}'. "
            "Use HEADING: markers for section titles with prose below."
        )
    elif intent == "txt":
        task_prefix = (
            f"Write clear, well-organised notes about '{topic}'. "
            "Plain text with clear headings."
        )
    elif intent == "summarise":
        task_prefix = (
            "You have been given the full text content of a workspace document. "
            "Please provide a clear, well-structured summary. "
            "Identify the main topics, key points, and any important details. "
            "Organise your summary with headings if the document has multiple sections."
        )
    else:
        task_prefix = ""

    full_query = f"{task_prefix}\n\n{query}".strip() if task_prefix else query
    msgs = [{"role": "system", "content": system_content}]
    msgs.extend(history[-6:])
    msgs.append({"role": "user", "content": full_query})
    return msgs


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 5 — Slide Parser
# ══════════════════════════════════════════════════════════════════════════════

def _strip_preamble(text: str) -> str:
    """Strip chatty LLM opener before the first structured content line."""
    trigger = re.compile(
        r"^(SLIDE\s*\d*\s*[:\-–]|#{1,3}\s|\*{2}|\-\s|\*\s|[1-9]\d*[\.\)]|CONTENT\s*[:\-–])",
        re.I | re.M
    )
    m = trigger.search(text)
    if m:
        return text[m.start():]
    return text


def _emergency_slides(topic: str, slide_count: int) -> list:
    """Generate topic-aware slide content when LLM output cannot be parsed."""
    t    = topic.strip().title() if topic else "This Topic"
    noun = t

    template = [
        (
            f"{t}: An Overview",
            [
                f"Introduction to {noun}",
                f"Why {noun} matters today",
                f"Key principles and scope",
                f"What this presentation covers",
            ],
        ),
        (
            f"Understanding {t}",
            [
                f"Core definition and background of {noun}",
                f"Historical development and evolution",
                f"Fundamental components explained",
                f"How {noun} works in practice",
            ],
        ),
        (
            f"Key Benefits of {t}",
            [
                f"Primary advantages of {noun}",
                f"Real-world impact and outcomes",
                f"Efficiency and effectiveness gains",
                f"Why organisations are adopting {noun}",
            ],
        ),
        (
            f"Applications of {t}",
            [
                f"Major industry use cases",
                f"Practical day-to-day implementations",
                f"Success stories and case studies",
                f"Emerging opportunities in {noun}",
            ],
        ),
        (
            f"Challenges & Considerations",
            [
                f"Common obstacles when applying {noun}",
                f"Risk factors to manage",
                f"Implementation and adoption hurdles",
                f"Strategies to overcome limitations",
            ],
        ),
        (
            f"The Future of {t}",
            [
                f"Emerging trends shaping {noun}",
                f"Predictions for the next five years",
                f"New opportunities on the horizon",
                f"How to prepare for what's next",
            ],
        ),
        (
            f"Key Takeaways",
            [
                f"{noun} is transforming the landscape",
                f"Summary of concepts covered today",
                f"Recommended next steps and actions",
                f"Thank you — questions welcome",
            ],
        ),
    ]

    result = template[:slide_count]
    idx = len(result) + 1
    while len(result) < slide_count:
        result.append((
            f"{t} — Insights Part {idx}",
            [
                f"Additional analysis of {noun}",
                f"Deeper dive into key concepts",
                f"Further considerations and context",
                f"Continued exploration of {noun}",
            ],
        ))
        idx += 1
    return result[:slide_count]


def parse_slides(llm_text: str, topic: str, slide_count: int) -> list:
    """Parse LLM output into (title, [bullets]) tuples."""
    clean_text = _strip_preamble(llm_text)
    lines      = clean_text.splitlines()

    # ── Pass 1: strict SLIDE/CONTENT format ──────────────────────────────────
    slides_data = []
    cur_title   = ""
    cur_bullets = []
    found_strict = False

    for line in lines:
        line = line.strip()
        m_slide = re.match(r"SLIDE\s*\d*\s*[:\-–]\s*(.+)", line, re.I)
        m_cont  = re.match(r"CONTENT\s*[:\-–]\s*(.+)",      line, re.I)
        if m_slide:
            if cur_title:
                slides_data.append((cur_title, cur_bullets))
            cur_title    = m_slide.group(1).strip().strip("*_")
            cur_bullets  = []
            found_strict = True
        elif m_cont and found_strict:
            cur_bullets = [b.strip() for b in m_cont.group(1).split("|") if b.strip()]

    if cur_title:
        slides_data.append((cur_title, cur_bullets))

    if slides_data and _has_real_content(slides_data):
        return _pad_slides(slides_data, topic, slide_count)

    # ── Pass 2: markdown/bold headings + following bullets ────────────────────
    slides_data  = []
    cur_title    = ""
    cur_bullets  = []

    heading_re  = re.compile(
        r"^(?:#{1,3}\s*|\*{1,2}[Ss]lide\s*\d*[:\-–]?\s*|\*{1,2}Slide\s*\d*[:\-–]?\s*)(.+?)(?:\*{1,2})?$"
    )
    bullet_re   = re.compile(r"^[\-\*\u2022\u25b8•]\s+(.+)")
    numbered_re = re.compile(r"^\d+[\.\)]\s+(.+)")

    for line in lines:
        line = line.strip()
        if not line:
            continue
        mh = heading_re.match(line)
        if mh and "|" not in line and len(mh.group(1)) < 80:
            ht = mh.group(1).strip().strip("*_:#")
            if cur_title:
                slides_data.append((cur_title, cur_bullets))
            cur_title   = ht
            cur_bullets = []
        else:
            mb = bullet_re.match(line) or numbered_re.match(line)
            if mb:
                if not cur_title:
                    cur_title = topic
                cur_bullets.append(mb.group(1).strip())
            elif "|" in line and cur_title:
                cur_bullets = [b.strip() for b in line.split("|") if b.strip()]

    if cur_title:
        slides_data.append((cur_title, cur_bullets))

    if slides_data and _has_real_content(slides_data):
        return _pad_slides(slides_data, topic, slide_count)

    # ── Pass 3: paragraph chunking ────────────────────────────────────────────
    paragraphs  = [p.strip() for p in re.split(r"\n{2,}", clean_text) if p.strip()]
    slides_data = []
    for i, para in enumerate(paragraphs[:slide_count]):
        sents  = re.split(r"(?<=[.!?])\s+", para)
        title  = sents[0][:60] if sents else f"Slide {i + 1}"
        buls   = [s[:120] for s in sents[1:5]] if len(sents) > 1 else [para[:120]]
        slides_data.append((title, buls))

    if slides_data and _has_real_content(slides_data):
        return _pad_slides(slides_data, topic, slide_count)

    # ── Emergency: topic-aware generated content ──────────────────────────────
    print_info("LLM output did not match any format — using topic-aware emergency content.")
    return _emergency_slides(topic, slide_count)


def _has_real_content(slides: list) -> bool:
    good = 0
    for _title, bullets in slides:
        if any(len(b.strip()) > 15 for b in bullets):
            good += 1
    return good >= max(1, len(slides) // 2)


def _pad_slides(slides: list, topic: str, target: int) -> list:
    if slides and not slides[0][1]:
        slides[0] = (slides[0][0], [topic, datetime.now().strftime("%B %Y")])
    _filler = re.compile(r"^\s*(sure|happy|glad|here|let me|i'll|of course|certainly)[^a-z]*$", re.I)
    if slides and _filler.match(slides[0][0]):
        slides[0] = (topic.title(), slides[0][1])
    if len(slides) < target:
        emergency = _emergency_slides(topic, target)
        while len(slides) < target:
            slides.append(emergency[len(slides)])
    return slides[:target]


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 6 — Output Generators
# ══════════════════════════════════════════════════════════════════════════════

def _safe_filename(topic: str) -> str:
    return re.sub(r"[^\w\s-]", "", topic).strip().replace(" ", "_")[:50]

def _timestamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


# ══════════════════════════════════════════════════════════════════════════════
#  6a — COM automation: write directly INTO PowerPoint on the user's PC
# ══════════════════════════════════════════════════════════════════════════════

_PP_LAYOUT_TITLE   = 1
_PP_LAYOUT_CONTENT = 2

def _rgb_long(r, g, b):
    return r + (g * 256) + (b * 65536)

# ── Colour palette ─────────────────────────────────────────────────────────────
_COL_BG      = _rgb_long(0x12, 0x12, 0x1E)   # dark navy background
_COL_ACCENT  = _rgb_long(0x6C, 0x63, 0xFF)   # purple accent
_COL_LIGHT   = _rgb_long(0xF0, 0xF0, 0xFF)   # near-white (cover title only)
_COL_SUB     = _rgb_long(0xA0, 0x9E, 0xC5)   # muted purple-grey (subtitles/footers)

# FIX: Slide titles use a vivid cyan-green so they stand out on the dark background
# and are never confused with the bullet text.
_COL_SLIDE_TITLE = _rgb_long(0x00, 0xE5, 0xCC)   # bright teal / cyan-green

# FIX: Bullet body text is now plain black for maximum readability.
_COL_BULLET  = _rgb_long(0x00, 0x00, 0x00)   # pure black


def _com_solid_fill(shape, color_long: int):
    """Apply a solid fill colour to a COM shape."""
    shape.Fill.Solid()
    shape.Fill.ForeColor.RGB = color_long


def _com_write_bullets_to_textbox(slide, bullets: list, left_in, top_in, width_in, height_in):
    """
    Add a white-background textbox and write bullets into it so the black text
    is always readable against the slide's dark background.
    """
    bx = slide.Shapes.AddTextbox(
        1,                          # msoTextOrientationHorizontal
        int(left_in   * 72),
        int(top_in    * 72),
        int(width_in  * 72),
        int(height_in * 72),
    )
    tf = bx.TextFrame
    tf.WordWrap = True

    # Give the textbox a white background so black text is always readable
    bx.Fill.Solid()
    bx.Fill.ForeColor.RGB = _rgb_long(0xFF, 0xFF, 0xFF)   # white background
    bx.Line.Visible = False

    bullet_text = "\r".join(f"\u25b8  {b}" for b in bullets)
    tf.TextRange.Text = bullet_text

    for i in range(1, tf.TextRange.Paragraphs().Count + 1):
        para = tf.TextRange.Paragraphs(i)
        para.Font.Size      = 15
        para.Font.Bold      = False
        para.Font.Color.RGB = _COL_BULLET   # BLACK
    return bx


def create_pptx_via_com(topic: str, llm_text: str, images: list, slide_count: int):
    """
    Open PowerPoint via COM automation, build the presentation INSIDE the running
    PowerPoint application, save the file, and return (path, app, prs).
    """
    if not WIN32_OK:
        raise RuntimeError(
            "pywin32 is not installed.\n"
            "Install it with:  pip install pywin32\n"
            "Then run:         python -m pywin32_postinstall -install"
        )

    slides_data = parse_slides(llm_text, topic, slide_count)
    fname       = OUTPUTS / f"{_safe_filename(topic)}_{_timestamp()}.pptx"

    pythoncom.CoInitialize()

    print_info("Launching PowerPoint via COM automation ...")
    pptApp             = win32com.client.Dispatch("PowerPoint.Application")
    pptApp.Visible     = True
    pptApp.WindowState = 1   # ppWindowNormal

    prs = pptApp.Presentations.Add()
    prs.PageSetup.SlideWidth  = 13.33 * 72
    prs.PageSetup.SlideHeight = 7.5   * 72

    img_iter = iter(images)

    for idx, (title, bullets) in enumerate(slides_data):
        print_info(f"  Writing slide {idx + 1}/{len(slides_data)}: {title[:55]}")

        total  = len(slides_data)
        BLANK_LAYOUT = 12
        slide = prs.Slides.Add(idx + 1, BLANK_LAYOUT)

        # ── Background ────────────────────────────────────────────────────────
        slide.Background.Fill.Solid()
        slide.Background.Fill.ForeColor.RGB = _COL_BG

        if idx == 0:
            # ══ COVER SLIDE ══════════════════════════════════════════════════

            # Branding label (top-left)
            brand = slide.Shapes.AddTextbox(1,
                int(0.4 * 72), int(0.25 * 72),
                int(9.0 * 72), int(0.45 * 72))
            brand.TextFrame.TextRange.Text = "PaperReady  |  Designed by Bhargav"
            brand.TextFrame.TextRange.Font.Size      = 10
            brand.TextFrame.TextRange.Font.Color.RGB = _COL_ACCENT

            # Accent bar (decorative, near bottom of cover)
            bar = slide.Shapes.AddShape(1,
                int(3.5 * 72), int(5.0 * 72),
                int(6.3 * 72), 4)
            _com_solid_fill(bar, _COL_ACCENT)
            bar.Line.Visible = False

            # FIX: Main title — large, centred, near-white on dark background
            title_box = slide.Shapes.AddTextbox(1,
                int(1.0 * 72), int(1.8 * 72),
                int(11.3 * 72), int(2.4 * 72))
            title_box.TextFrame.WordWrap = True
            title_box.TextFrame.TextRange.Text = title          # <-- title IS written here
            title_box.TextFrame.TextRange.Font.Size      = 40
            title_box.TextFrame.TextRange.Font.Bold      = True
            title_box.TextFrame.TextRange.Font.Color.RGB = _COL_LIGHT   # near-white on cover
            title_box.TextFrame.TextRange.ParagraphFormat.Alignment = 2  # ppAlignCenter

            # Subtitle / date
            sub_box = slide.Shapes.AddTextbox(1,
                int(1.0 * 72), int(4.4 * 72),
                int(11.3 * 72), int(0.8 * 72))
            sub_box.TextFrame.TextRange.Text = f"{topic}  |  {datetime.now().strftime('%B %d, %Y')}"
            sub_box.TextFrame.TextRange.Font.Size      = 16
            sub_box.TextFrame.TextRange.Font.Color.RGB = _COL_SUB
            sub_box.TextFrame.TextRange.ParagraphFormat.Alignment = 2  # ppAlignCenter

        else:
            # ══ CONTENT SLIDE ════════════════════════════════════════════════

            # White background panel behind the title so black/teal text is clear
            title_bg = slide.Shapes.AddShape(1,
                0, 0, int(13.33 * 72), int(1.1 * 72))
            _com_solid_fill(title_bg, _rgb_long(0xFF, 0xFF, 0xFF))   # white banner
            title_bg.Line.Visible = False

            # Accent bar (top edge over the white banner)
            top_bar = slide.Shapes.AddShape(1, 0, 0, int(13.33 * 72), 6)
            _com_solid_fill(top_bar, _COL_ACCENT)
            top_bar.Line.Visible = False

            # FIX: Slide title — written explicitly, teal colour, on white banner
            title_box = slide.Shapes.AddTextbox(1,
                int(0.5 * 72), int(0.05 * 72),
                int(11.5 * 72), int(1.0 * 72))
            title_box.TextFrame.WordWrap = True
            title_box.TextFrame.TextRange.Text = title          # <-- title IS written here
            title_box.TextFrame.TextRange.Font.Size      = 26
            title_box.TextFrame.TextRange.Font.Bold      = True
            title_box.TextFrame.TextRange.Font.Color.RGB = _COL_SLIDE_TITLE  # teal

            # Slide counter (top-right, inside white banner)
            ctr = slide.Shapes.AddTextbox(1,
                int(12.0 * 72), int(0.15 * 72),
                int(1.1  * 72), int(0.45 * 72))
            ctr.TextFrame.TextRange.Text = f"{idx} / {total - 1}"
            ctr.TextFrame.TextRange.Font.Size      = 11
            ctr.TextFrame.TextRange.Font.Color.RGB = _COL_ACCENT

            # Footer (dark strip at the very bottom)
            footer = slide.Shapes.AddTextbox(1,
                int(0.4 * 72), int(7.1 * 72),
                int(12.0 * 72), int(0.35 * 72))
            footer.TextFrame.TextRange.Text = (
                "PaperReady  |  Designed by Bhargav  |  Phi-3 Mini (Microsoft)"
            )
            footer.TextFrame.TextRange.Font.Size      = 9
            footer.TextFrame.TextRange.Font.Color.RGB = _COL_SUB

            # Determine content width (leave room for image if present)
            img_path = next(img_iter, None)
            content_width_in = 7.8 if (img_path and img_path.exists()) else 12.3

            # ── BULLET CONTENT — white background, black text ──────────────
            if bullets:
                _com_write_bullets_to_textbox(
                    slide, bullets,
                    left_in=0.5, top_in=1.2,
                    width_in=content_width_in, height_in=5.6
                )
            else:
                no_content = slide.Shapes.AddTextbox(1,
                    int(0.5 * 72), int(1.2 * 72),
                    int(content_width_in * 72), int(5.6 * 72))
                no_content.TextFrame.TextRange.Text = "(No content generated for this slide.)"
                no_content.TextFrame.TextRange.Font.Size      = 15
                no_content.TextFrame.TextRange.Font.Color.RGB = _COL_SUB

            # ── Image (right panel) ──────────────────────────────────────────
            if img_path and img_path.exists():
                try:
                    slide.Shapes.AddPicture(
                        str(img_path.resolve()),
                        LinkToFile=False,
                        SaveWithDocument=True,
                        Left=int(8.8 * 72), Top=int(1.25 * 72),
                        Width=int(4.0 * 72), Height=int(5.5 * 72)
                    )
                    print_info(f"    Image inserted: {img_path.name}")
                except Exception as e:
                    print_err(f"    Image failed ({img_path.name}): {e}")

    prs.SaveAs(str(fname.resolve()))
    print_ok(f"PowerPoint saved via COM  ->  {fname.name}")
    return fname, pptApp, prs


# ══════════════════════════════════════════════════════════════════════════════
#  6b — COM automation: write directly INTO Word on the user's PC
# ══════════════════════════════════════════════════════════════════════════════

def create_docx_via_com(topic: str, llm_text: str, images: list):
    if not WIN32_OK:
        raise RuntimeError(
            "pywin32 is not installed.\n"
            "Install it with:  pip install pywin32\n"
            "Then run:         python -m pywin32_postinstall -install"
        )

    fname = OUTPUTS / f"{_safe_filename(topic)}_{_timestamp()}.docx"

    pythoncom.CoInitialize()
    print_info("Launching Microsoft Word via COM automation ...")

    wordApp             = win32com.client.Dispatch("Word.Application")
    wordApp.Visible     = True
    wordApp.WindowState = 1

    doc = wordApp.Documents.Add()

    def _append_para(text: str, style_name: str = "Normal",
                     size: int = 11, bold: bool = False, color_long: int = 0x000000):
        rng = doc.Content
        rng.Collapse(0)
        rng.InsertParagraphAfter()
        rng.Collapse(0)
        rng.Text = text
        rng.Font.Size  = size
        rng.Font.Bold  = bold
        rng.Font.Color = color_long
        try:
            rng.Style = doc.Styles(style_name)
        except Exception:
            pass

    first_para = doc.Paragraphs(1).Range
    first_para.Text        = topic
    first_para.Font.Size   = 26
    first_para.Font.Bold   = True
    first_para.Font.Color  = _COL_ACCENT
    try:
        first_para.Style = doc.Styles("Title")
    except Exception:
        pass

    _append_para(
        f"Generated by PaperReady  |  Designed by Bhargav  |  "
        f"Phi-3 Mini (Microsoft)  |  {datetime.now().strftime('%d %b %Y %H:%M')}",
        size=9,
        color_long=_rgb_long(0x88, 0x88, 0xAA)
    )
    _append_para("")

    heading_re  = re.compile(r"HEADING\s*[:\-–]\s*(.+)", re.I)
    img_iter    = iter(images)
    first_image = True

    for raw_line in llm_text.splitlines():
        line = raw_line.strip()
        if not line:
            _append_para("")
            continue

        m = heading_re.match(line)
        if m:
            heading_text = m.group(1).strip()
            _append_para(heading_text,
                         style_name="Heading 1",
                         size=14,
                         bold=True,
                         color_long=_COL_ACCENT)

            if first_image:
                img_path = next(img_iter, None)
                if img_path and img_path.exists():
                    try:
                        rng = doc.Content
                        rng.Collapse(0)
                        rng.InsertParagraphAfter()
                        rng.Collapse(0)
                        rng.InlineShapes.AddPicture(
                            FileName=str(img_path.resolve()),
                            LinkToFile=False,
                            SaveWithDocument=True
                        )
                        _append_para("")
                        first_image = False
                        print_info(f"    Image inserted in Word: {img_path.name}")
                    except Exception as e:
                        print_err(f"    Word image failed: {e}")
        else:
            _append_para(line)

    for img_path in img_iter:
        if img_path.exists():
            try:
                rng = doc.Content
                rng.Collapse(0)
                rng.InsertParagraphAfter()
                rng.Collapse(0)
                rng.InlineShapes.AddPicture(
                    FileName=str(img_path.resolve()),
                    LinkToFile=False,
                    SaveWithDocument=True
                )
            except Exception:
                pass

    doc.SaveAs2(str(fname.resolve()))
    print_ok(f"Word document saved via COM  ->  {fname.name}")
    return fname, wordApp, doc


# ══════════════════════════════════════════════════════════════════════════════
#  6c — python-pptx fallback (when COM / win32com is unavailable)
# ══════════════════════════════════════════════════════════════════════════════

def create_pptx_fallback(topic: str, llm_text: str, images: list,
                         slide_count: int) -> Path:
    if not PPTX_OK:
        raise RuntimeError("python-pptx not installed.  Run:  pip install python-pptx")

    slides_data = parse_slides(llm_text, topic, slide_count)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK_BG      = RGBColor(0x12, 0x12, 0x1E)
    ACCENT       = RGBColor(0x6C, 0x63, 0xFF)
    LIGHT_TXT    = RGBColor(0xF0, 0xF0, 0xFF)
    SUB_TXT      = RGBColor(0xA0, 0x9E, 0xC5)
    # FIX: teal title colour and black bullet colour
    TITLE_CLR    = RGBColor(0x00, 0xE5, 0xCC)   # bright teal
    BULLET_CLR   = RGBColor(0x00, 0x00, 0x00)   # pure black
    WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

    blank_layout = prs.slide_layouts[6]

    def hex_bg(slide, color):
        bg = slide.background; fill = bg.fill
        fill.solid(); fill.fore_color.rgb = color

    def add_tb(slide, text, l, t, w, h, fs=18, bold=False,
               color=None, align=PP_ALIGN.LEFT, bg_color=None):
        if color is None: color = LIGHT_TXT
        txBox = slide.shapes.add_textbox(l, t, w, h)
        tf = txBox.text_frame; tf.word_wrap = True
        if bg_color is not None:
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = bg_color
            txBox.line.fill.background()
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run()
        run.text = text; run.font.size = Pt(fs)
        run.font.bold = bold; run.font.color.rgb = color
        return txBox

    def add_rect(slide, l, t, w, h, color):
        shape = slide.shapes.add_shape(1, l, t, w, h)
        shape.fill.solid(); shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_bullets(slide, bullets, l, t, w, h):
        txBox = slide.shapes.add_textbox(l, t, w, h)
        tf = txBox.text_frame; tf.word_wrap = True
        # White background so black text is readable on the dark slide
        txBox.fill.solid(); txBox.fill.fore_color.rgb = WHITE
        txBox.line.fill.background()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = f"  \u25b8  {b}"
            run.font.size = Pt(15)
            run.font.color.rgb = BULLET_CLR   # BLACK
        return txBox

    img_iter = iter(images)
    total    = len(slides_data)

    for idx, (title, bullets) in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        hex_bg(slide, DARK_BG)

        if idx == 0:
            # Cover slide
            add_tb(slide, "PaperReady  |  Designed by Bhargav",
                   Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                   fs=11, color=ACCENT)
            # FIX: title is explicitly written here
            add_tb(slide, title,
                   Inches(1), Inches(2.2), Inches(11.3), Inches(2.2),
                   fs=44, bold=True, color=LIGHT_TXT, align=PP_ALIGN.CENTER)
            add_tb(slide, f"{topic}\n{datetime.now().strftime('%B %d, %Y')}",
                   Inches(1), Inches(4.6), Inches(11.3), Inches(1.0),
                   fs=18, color=SUB_TXT, align=PP_ALIGN.CENTER)
            add_rect(slide, Inches(3.5), Inches(5.0), Inches(6.3), Pt(3), ACCENT)
        else:
            # Content slide — white banner at top for the title
            add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), WHITE)
            # Thin accent bar along very top edge
            add_rect(slide, Inches(0), Inches(0), Inches(13.33), Pt(5), ACCENT)

            # FIX: slide title — teal on white banner
            add_tb(slide, title,
                   Inches(0.5), Inches(0.05), Inches(11.5), Inches(1.0),
                   fs=26, bold=True, color=TITLE_CLR)

            add_tb(slide, f"{idx} / {total - 1}",
                   Inches(12.0), Inches(0.2), Inches(1.1), Inches(0.5),
                   fs=11, color=ACCENT, align=PP_ALIGN.RIGHT)
            add_tb(slide, "PaperReady  |  Designed by Bhargav  |  Phi-3 Mini",
                   Inches(0.5), Inches(7.1), Inches(12), Inches(0.35),
                   fs=9, color=SUB_TXT)

            img_path      = next(img_iter, None)
            content_width = Inches(8.0) if (img_path and img_path.exists()) else Inches(12.3)

            # FIX: bullets with white bg + black text
            if bullets:
                add_bullets(slide, bullets, Inches(0.5), Inches(1.3), content_width, Inches(5.4))
            else:
                add_tb(slide, "No content generated for this slide.",
                       Inches(0.5), Inches(1.3), content_width, Inches(5.4),
                       fs=15, color=SUB_TXT)

            if img_path and img_path.exists():
                try:
                    slide.shapes.add_picture(
                        str(img_path), Inches(8.8), Inches(1.25), Inches(4.0), Inches(5.5))
                    print_info(f"Image inserted slide {idx}: {img_path.name}")
                except Exception as e:
                    print_err(f"Image insert failed: {e}")

    fname = OUTPUTS / f"{_safe_filename(topic)}_{_timestamp()}.pptx"
    prs.save(str(fname))
    return fname


# ══════════════════════════════════════════════════════════════════════════════
#  6d — python-docx fallback
# ══════════════════════════════════════════════════════════════════════════════

def create_docx_fallback(topic: str, llm_text: str, images: list) -> Path:
    if not DOCX_OK:
        raise RuntimeError("python-docx not installed.  Run:  pip install python-docx")

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"; style.font.size = DocPt(11)

    tp = doc.add_heading(topic, level=0)
    tp.runs[0].font.color.rgb = DocRGB(0x6C, 0x63, 0xFF)

    meta = doc.add_paragraph(
        f"Generated by PaperReady  |  Designed by Bhargav  |  "
        f"Phi-3 Mini (Microsoft)  |  {datetime.now().strftime('%d %b %Y %H:%M')}"
    )
    meta.runs[0].font.color.rgb = DocRGB(0x88, 0x88, 0xAA)
    meta.runs[0].font.size      = DocPt(9)
    doc.add_paragraph("")

    img_iter          = iter(images)
    heading_re        = re.compile(r"HEADING\s*[:\-–]\s*(.+)", re.I)
    first_image_added = False

    for line in llm_text.splitlines():
        ls = line.strip()
        if not ls:
            doc.add_paragraph(""); continue
        m = heading_re.match(ls)
        if m:
            h = doc.add_heading(m.group(1).strip(), level=1)
            if h.runs: h.runs[0].font.color.rgb = DocRGB(0x6C, 0x63, 0xFF)
            if not first_image_added:
                img_path = next(img_iter, None)
                if img_path and img_path.exists():
                    try:
                        doc.add_picture(str(img_path), width=DocInches(5))
                        doc.add_paragraph(""); first_image_added = True
                        print_info(f"Image in document: {img_path.name}")
                    except Exception as e:
                        print_err(f"Image failed: {e}")
        else:
            p = doc.add_paragraph(ls); p.paragraph_format.space_after = DocPt(6)

    for img_path in img_iter:
        if img_path.exists():
            try:
                doc.add_paragraph("")
                doc.add_picture(str(img_path), width=DocInches(5))
            except Exception:
                pass

    fname = OUTPUTS / f"{_safe_filename(topic)}_{_timestamp()}.docx"
    doc.save(str(fname))
    return fname


# ══════════════════════════════════════════════════════════════════════════════
#  6e — TXT writer
# ══════════════════════════════════════════════════════════════════════════════

def create_txt(topic: str, llm_text: str) -> Path:
    rule   = "=" * (W - LEFT)
    header = (
        f"{rule}\n"
        f"  {topic}\n"
        f"  Generated : {datetime.now().strftime('%d %b %Y %H:%M')}\n"
        f"  PaperReady Editor  |  Designed by Bhargav  |  Phi-3 Mini (Microsoft)\n"
        f"{rule}\n\n"
    )
    content = header + llm_text.strip() + "\n"
    fname   = OUTPUTS / f"{_safe_filename(topic)}_{_timestamp()}.txt"
    fname.write_text(content, encoding="utf-8")
    return fname


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 7 — About / Identity Response
# ══════════════════════════════════════════════════════════════════════════════

def print_about():
    print()
    print(_hr("="))
    print(_center_line(f"  About {IDENTITY['name']}  "))
    print(_hr("="))
    print(_box_line())
    print(_center_line(f"Designed & Built by  :  {IDENTITY['designer']}"))
    print(_center_line(f"Base Language Model  :  {IDENTITY['base_model']}"))
    print(_center_line(f"Version              :  {IDENTITY['version']}"))
    print(_box_line())
    print(_box_line())

    for wrapped_line in _box_wrap_lines(IDENTITY["purpose"]):
        print(wrapped_line)

    print(_box_line())
    print(_box_line("  What I can do for you:"))

    for cap in IDENTITY["capabilities"]:
        prefix  = "    ->  "
        avail   = INNER - len(prefix)
        words   = cap.split()
        line    = ""
        first   = True
        for word in words:
            candidate = f"{line} {word}".strip() if line else word
            if len(candidate) > avail:
                if first:
                    print(_box_line(prefix + line))
                    first = False
                else:
                    print(_box_line((" " * len(prefix)) + line))
                line = word
            else:
                line = candidate
        tail = (prefix if first else " " * len(prefix)) + line
        if tail.strip():
            print(_box_line(tail))

    print(_box_line())
    print(_hr("="))
    print()

    if WIN32_OK:
        print_ok("COM automation (pywin32) detected — PPT/Word will be written LIVE in the app.")
    else:
        print_info("pywin32 not found — falling back to python-pptx / python-docx file creation.")
        print_info("Install with:  pip install pywin32  then  python -m pywin32_postinstall -install")
    print()


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 8 — LLM Loader
# ══════════════════════════════════════════════════════════════════════════════

def load_model():
    if not LLAMA_OK:
        raise RuntimeError("llama_cpp not installed.  Run:  pip install llama-cpp-python")
    print_info("Loading Phi-3 Mini Q4 model (~2.2 GB) from Hugging Face cache ...")
    print_info("First run will download; subsequent runs use local cache.")
    llm = Llama.from_pretrained(
        repo_id   = "microsoft/Phi-3-mini-4k-instruct-gguf",
        filename  = "Phi-3-mini-4k-instruct-q4.gguf",
        n_ctx     = 4096,
        n_threads = min(8, psutil.cpu_count(logical=False) or 4),
        verbose   = False,
    )
    return llm


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 9 — Banner
# ══════════════════════════════════════════════════════════════════════════════

def print_banner():
    print()
    print(_hr("="))
    print()
    print(_figlet_banner("PaperReady", font="slant"))
    subtitle = "RAG-Powered Local AI  |  Phi-3 Mini  |  Designed by Bhargav"
    print(PAD + subtitle.center(W - LEFT))
    print()
    print(_hr("="))
    print(_box_line())
    print(_center_line("COMMANDS"))
    print(_box_line())
    print(_box_line("Type any question                 ->  chat with the AI"))
    print(_box_line("create a 3 slide ppt about X      ->  opens PPT, writes 3 slides"))
    print(_box_line("write a word document about X     ->  opens Word, writes document"))
    print(_box_line("write a text file about X         ->  saves .txt"))
    print(_box_line("summarize report.pptx             ->  reads & summarises workspace file"))
    print(_box_line("summarize notes.docx              ->  reads & summarises workspace file"))
    print(_box_line("who are you / what can you do     ->  about PaperReady"))
    print(_box_line("reload  ->  re-scan workspace/"))
    print(_box_line("exit / quit  ->  close PaperReady"))
    print(_box_line())

    com_status = "COM (live in-app)" if WIN32_OK else "python-pptx / python-docx (file only)"
    for wl in _box_wrap_lines(f"  PPT/DOCX mode: {com_status}", indent="  "):
        print(wl)

    print(_box_line())
    print(_hr("="))
    print()


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 10 — Main Loop
# ══════════════════════════════════════════════════════════════════════════════

def run_llm_streaming(llm, messages: list) -> str:
    print()
    print(f"{PAD}  [ AI ]")
    print(f"{PAD}  ", end="", flush=True)

    chunks  = llm.create_chat_completion(
        messages    = messages,
        stream      = True,
        max_tokens  = 2048,
        temperature = 0.4,
    )
    full    = ""
    col     = LEFT + 2
    wrap_at = W - 2

    for chunk in chunks:
        delta = chunk["choices"][0]["delta"]
        if "content" in delta:
            piece = delta["content"]
            for char in piece:
                if char == "\n":
                    print(); print(f"{PAD}  ", end="", flush=True); col = LEFT + 2
                else:
                    print(char, end="", flush=True); col += 1
                    if col >= wrap_at:
                        print(); print(f"{PAD}  ", end="", flush=True); col = LEFT + 2
            full += piece

    print(); print()
    return full


def prompt_user() -> str:
    print(f"{PAD}  [ YOU ]")
    return input(f"{PAD}  ").strip()


def handle_file_output(intent: str, fname: Path, com_app=None, com_doc=None):
    app_names = {"pptx": "Microsoft PowerPoint",
                 "docx": "Microsoft Word",
                 "txt" : "Notepad / text editor"}
    app_name  = app_names.get(intent, "the default application")

    print_ok(f"File saved  ->  {fname}")
    print_info(f"File size   ->  {fname.stat().st_size // 1024} KB")

    if com_app is not None:
        print_ok(f"{app_name} is open with your content.")
        print_info("The file has been saved. You can continue editing in the app.")
    else:
        if ask_permission(f"Open '{fname.name}' in {app_name}?"):
            if open_file_in_app(fname, app_hint=intent):
                print_ok(f"Opened in {app_name}.")
            else:
                print_err("Auto-open failed. Open manually from outputs/")
        else:
            print_info(f"File is ready at:  {fname}")
            print_info("Open it manually from the outputs/ folder whenever you like.")


# ══════════════════════════════════════════════════════════════════════════════
# NEW: handle_summarise_request
# ══════════════════════════════════════════════════════════════════════════════

def handle_summarise_request(llm, rag: WorkspaceRAG, user_input: str, history: list):
    """
    Detect the target file from user_input, extract its text, feed it to the LLM,
    and print a summary.  Returns the updated history list.
    """
    filename_hint = detect_summarise_request(user_input)
    if not filename_hint:
        print_err("Could not identify a filename in your request.")
        print_info("Try: 'summarize report.pptx' or 'read notes.docx'")
        return history

    target_path = rag.find_file(filename_hint)
    if target_path is None:
        print_err(f"File not found in workspace/:  {filename_hint}")
        print_info("Place the file in the workspace/ folder and type 'reload', then try again.")
        return history

    print_section(f"Reading file  ->  {target_path.name}")
    extracted_text = extract_workspace_file(target_path)

    if not extracted_text.strip():
        print_err(f"Could not extract any text from {target_path.name}.")
        return history

    word_count = len(extracted_text.split())
    print_info(f"Extracted {word_count} words from {target_path.name}")

    # Truncate to fit in the context window safely (keep ~3000 words)
    MAX_WORDS = 3000
    if word_count > MAX_WORDS:
        truncated = " ".join(extracted_text.split()[:MAX_WORDS])
        print_info(f"(Truncated to {MAX_WORDS} words to fit LLM context window)")
    else:
        truncated = extracted_text

    summarise_query = (
        f"Here is the full content of the file '{target_path.name}':\n\n"
        f"{truncated}\n\n"
        f"Please provide a clear, well-structured summary of this document. "
        f"Identify the main topics, key points, and any important details."
    )

    messages = build_messages(
        history, summarise_query, context="", intent="summarise",
        topic=target_path.stem
    )

    try:
        llm_response = run_llm_streaming(llm, messages)
    except Exception as e:
        print_err(f"Inference error: {e}")
        return history

    history.append({"role": "user",      "content": user_input})
    history.append({"role": "assistant", "content": llm_response})
    return history


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 11 — Main Loop
# ══════════════════════════════════════════════════════════════════════════════

def main():
    print()
    print(_hr("="))
    print(_center_line("PaperReady Editor  --  initialising ..."))
    print(_center_line("Designed by Bhargav  |  Phi-3 Mini (Microsoft)"))
    print(_hr("="))
    print()

    rag = WorkspaceRAG()
    print()

    try:
        llm = load_model()
        print_ok("Model loaded successfully!")
    except Exception as e:
        print_err(f"Model load failed: {e}")
        sys.exit(1)

    print_banner()

    history = []

    while True:
        try:
            user_input = prompt_user()
        except (EOFError, KeyboardInterrupt):
            print(); print_info("Goodbye! — PaperReady by Bhargav"); break

        if not user_input:
            continue

        if user_input.lower() in ("exit", "quit"):
            print_info("Goodbye! — PaperReady by Bhargav"); break

        if user_input.lower() == "reload":
            print_section("Reloading Workspace Index"); rag.reload(); continue

        if ABOUT_TRIGGERS.search(user_input):
            print_about()
            history.append({"role": "user",      "content": user_input})
            history.append({"role": "assistant", "content": IDENTITY["purpose"]})
            print(_hr("-")); continue

        intent      = detect_intent(user_input)
        topic       = extract_topic(user_input)
        slide_count = extract_slide_count(user_input) if intent == "pptx" else 6

        # ── NEW: file summarisation path ──────────────────────────────────────
        if intent == "summarise":
            history = handle_summarise_request(llm, rag, user_input, history)
            print(_hr("-"))
            continue

        if intent != "chat":
            print_section(
                f"Task  ->  {intent.upper()}  |  Topic: {topic}"
                + (f"  |  Slides: {slide_count}" if intent == "pptx" else "")
            )

        context = rag.retrieve_text(user_input)
        images  = rag.retrieve_images(topic) if intent in ("pptx", "docx") else []

        if images:
            print_info(f"Matched {len(images)} workspace image(s): " +
                       ", ".join(p.name for p in images))
        elif intent in ("pptx", "docx"):
            print_info("No matching images found in workspace/ for this topic.")

        messages = build_messages(history, user_input, context, intent, topic, slide_count)

        try:
            llm_response = run_llm_streaming(llm, messages)
        except Exception as e:
            print_err(f"Inference error: {e}"); continue

        history.append({"role": "user",      "content": user_input})
        history.append({"role": "assistant", "content": llm_response})

        if intent == "pptx":
            print_section(f"Building PowerPoint  ->  {topic}  ({slide_count} slides)")
            preview_lines = [l for l in llm_response.splitlines() if l.strip()][:8]
            print_info("LLM output preview (first 8 non-empty lines):")
            for pl in preview_lines:
                print(f"{PAD}    {pl[:80]}")
            print()
            com_app = com_doc = None
            try:
                if WIN32_OK:
                    fname, com_app, com_doc = create_pptx_via_com(
                        topic, llm_response, images, slide_count)
                else:
                    print_info("pywin32 not available — using python-pptx fallback.")
                    fname = create_pptx_fallback(topic, llm_response, images, slide_count)
                handle_file_output("pptx", fname, com_app, com_doc)
            except Exception as e:
                print_err(f"PPTX error: {e}")

        elif intent == "docx":
            print_section(f"Building Word Document  ->  {topic}")
            com_app = com_doc = None
            try:
                if WIN32_OK:
                    fname, com_app, com_doc = create_docx_via_com(
                        topic, llm_response, images)
                else:
                    print_info("pywin32 not available — using python-docx fallback.")
                    fname = create_docx_fallback(topic, llm_response, images)
                handle_file_output("docx", fname, com_app, com_doc)
            except Exception as e:
                print_err(f"DOCX error: {e}")

        elif intent == "txt":
            print_section(f"Writing Text File  ->  {topic}")
            try:
                fname = create_txt(topic, llm_response)
                handle_file_output("txt", fname)
            except Exception as e:
                print_err(f"TXT error: {e}")

        print(_hr("-"))


if __name__ == "__main__":
    main()