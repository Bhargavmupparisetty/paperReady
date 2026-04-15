from pathlib import Path
from paperready.utils import print_err

try:
    from pptx import Presentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    from docx import Document
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

def _extract_text_from_pptx(path: Path) -> str:
    if not PPTX_OK:
        return ""
    try:
        prs = Presentation(str(path))
        lines = [f"[PowerPoint: {path.name}]"]
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt:
                            lines.append(txt)
        return "\n".join(lines)
    except Exception as e:
        print_err(f"Could not read {path.name}: {e}")
        return ""

def _extract_text_from_docx(path: Path) -> str:
    if not DOCX_OK:
        return ""
    try:
        doc = Document(str(path))
        lines = [f"[Word Document: {path.name}]"]
        for para in doc.paragraphs:
            txt = para.text.strip()
            if txt:
                lines.append(txt)
        return "\n".join(lines)
    except Exception as e:
        print_err(f"Could not read {path.name}: {e}")
        return ""

def _extract_text_from_txt(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        print_err(f"Could not read {path.name}: {e}")
        return ""

def extract_workspace_file(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pptx":
        return _extract_text_from_pptx(path)
    elif ext == ".docx":
        return _extract_text_from_docx(path)
    else:
        return _extract_text_from_txt(path)
