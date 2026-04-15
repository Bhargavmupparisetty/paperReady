import re
from pathlib import Path
from datetime import datetime
from paperready.config import OUTPUTS, W, LEFT
from paperready.utils import print_info, print_ok, print_err
from paperready.parsers import parse_slides

try:
    import win32com.client
    import pythoncom
    WIN32_OK = True
except ImportError:
    WIN32_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    from docx import Document
    from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGB
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

def _safe_filename(topic: str) -> str:
    return re.sub(r"[^\w\s-]", "", topic).strip().replace(" ", "_")[:50]

def _timestamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")

def _rgb_long(r, g, b):
    return r + (g * 256) + (b * 65536)

_COL_BG = _rgb_long(0x12, 0x12, 0x1E)
_COL_ACCENT = _rgb_long(0x6C, 0x63, 0xFF)
_COL_LIGHT = _rgb_long(0xF0, 0xF0, 0xFF)
_COL_SUB = _rgb_long(0xA0, 0x9E, 0xC5)
_COL_SLIDE_TITLE = _rgb_long(0x00, 0xE5, 0xCC)
_COL_BULLET = _rgb_long(0x00, 0x00, 0x00)
_COL_COVER_TITLE = _rgb_long(0x29, 0xB6, 0xFF)

def _com_solid_fill(shape, color_long: int):
    shape.Fill.Solid()
    shape.Fill.ForeColor.RGB = color_long

def _com_write_bullets_to_textbox(slide, bullets: list, left_in, top_in, width_in, height_in):
    bx = slide.Shapes.AddTextbox(
        1,
        int(left_in * 72),
        int(top_in * 72),
        int(width_in * 72),
        int(height_in * 72),
    )
    tf = bx.TextFrame
    tf.WordWrap = True
    bx.Fill.Solid()
    bx.Fill.ForeColor.RGB = _rgb_long(0xFF, 0xFF, 0xFF)
    bx.Line.Visible = False
    bullet_text = "\r".join(f"\u25b8  {b}" for b in bullets)
    tf.TextRange.Text = bullet_text
    for i in range(1, tf.TextRange.Paragraphs().Count + 1):
        para = tf.TextRange.Paragraphs(i)
        para.Font.Size = 15
        para.Font.Bold = False
        para.Font.Color.RGB = _COL_BULLET
    return bx

def create_pptx_via_com(topic: str, llm_text: str, images: list, slide_count: int):
    if not WIN32_OK:
        raise RuntimeError("pywin32 is not installed.")
    slides_data = parse_slides(llm_text, topic, slide_count)
    fname = OUTPUTS / f"{_safe_filename(topic)}_{_timestamp()}.pptx"
    pythoncom.CoInitialize()
    print_info("Launching PowerPoint via COM automation ...")
    pptApp = win32com.client.Dispatch("PowerPoint.Application")
    pptApp.Visible = True
    pptApp.WindowState = 1
    prs = pptApp.Presentations.Add()
    prs.PageSetup.SlideWidth = 13.33 * 72
    prs.PageSetup.SlideHeight = 7.5 * 72
    img_iter = iter(images)
    for idx, (title, bullets) in enumerate(slides_data):
        print_info(f"  Writing slide {idx + 1}/{len(slides_data)}: {title[:55]}")
        total = len(slides_data)
        BLANK_LAYOUT = 12
        slide = prs.Slides.Add(idx + 1, BLANK_LAYOUT)
        slide.Background.Fill.Solid()
        slide.Background.Fill.ForeColor.RGB = _COL_BG
        if idx == 0:
            brand = slide.Shapes.AddTextbox(1, int(0.4 * 72), int(0.25 * 72), int(9.0 * 72), int(0.45 * 72))
            brand.TextFrame.TextRange.Text = "PaperReady  |  Designed by Bhargav"
            brand.TextFrame.TextRange.Font.Size = 10
            brand.TextFrame.TextRange.Font.Color.RGB = _COL_ACCENT
            bar = slide.Shapes.AddShape(1, int(3.5 * 72), int(5.0 * 72), int(6.3 * 72), 4)
            _com_solid_fill(bar, _COL_ACCENT)
            bar.Line.Visible = False
            title_box = slide.Shapes.AddTextbox(1, int(1.0 * 72), int(1.8 * 72), int(11.3 * 72), int(2.4 * 72))
            title_box.TextFrame.WordWrap = True
            title_box.TextFrame.TextRange.Text = topic.strip().title()
            title_box.TextFrame.TextRange.Font.Size = 40
            title_box.TextFrame.TextRange.Font.Bold = True
            title_box.TextFrame.TextRange.Font.Color.RGB = _COL_COVER_TITLE
            title_box.TextFrame.TextRange.ParagraphFormat.Alignment = 2
            sub_box = slide.Shapes.AddTextbox(1, int(1.0 * 72), int(4.4 * 72), int(11.3 * 72), int(0.8 * 72))
            sub_box.TextFrame.TextRange.Text = f"{topic.strip().title()}  |  {datetime.now().strftime('%B %d, %Y')}"
            sub_box.TextFrame.TextRange.Font.Size = 16
            sub_box.TextFrame.TextRange.Font.Color.RGB = _COL_SUB
            sub_box.TextFrame.TextRange.ParagraphFormat.Alignment = 2
        else:
            title_bg = slide.Shapes.AddShape(1, 0, 0, int(13.33 * 72), int(1.1 * 72))
            _com_solid_fill(title_bg, _rgb_long(0xFF, 0xFF, 0xFF))
            title_bg.Line.Visible = False
            top_bar = slide.Shapes.AddShape(1, 0, 0, int(13.33 * 72), 6)
            _com_solid_fill(top_bar, _COL_ACCENT)
            top_bar.Line.Visible = False
            title_box = slide.Shapes.AddTextbox(1, int(0.5 * 72), int(0.05 * 72), int(11.5 * 72), int(1.0 * 72))
            title_box.TextFrame.WordWrap = True
            title_box.TextFrame.TextRange.Text = title
            title_box.TextFrame.TextRange.Font.Size = 26
            title_box.TextFrame.TextRange.Font.Bold = True
            title_box.TextFrame.TextRange.Font.Color.RGB = _COL_SLIDE_TITLE
            ctr = slide.Shapes.AddTextbox(1, int(12.0 * 72), int(0.15 * 72), int(1.1 * 72), int(0.45 * 72))
            ctr.TextFrame.TextRange.Text = f"{idx} / {total - 1}"
            ctr.TextFrame.TextRange.Font.Size = 11
            ctr.TextFrame.TextRange.Font.Color.RGB = _COL_ACCENT
            footer = slide.Shapes.AddTextbox(1, int(0.4 * 72), int(7.1 * 72), int(12.0 * 72), int(0.35 * 72))
            footer.TextFrame.TextRange.Text = "PaperReady  |  Designed by Bhargav  |  Phi-3 Mini (Microsoft)"
            footer.TextFrame.TextRange.Font.Size = 9
            footer.TextFrame.TextRange.Font.Color.RGB = _COL_SUB
            img_path = next(img_iter, None)
            content_width_in = 7.8 if (img_path and img_path.exists()) else 12.3
            if bullets:
                _com_write_bullets_to_textbox(slide, bullets, left_in=0.5, top_in=1.2, width_in=content_width_in, height_in=5.6)
            else:
                no_content = slide.Shapes.AddTextbox(1, int(0.5 * 72), int(1.2 * 72), int(content_width_in * 72), int(5.6 * 72))
                no_content.TextFrame.TextRange.Text = "(No content generated for this slide.)"
                no_content.TextFrame.TextRange.Font.Size = 15
                no_content.TextFrame.TextRange.Font.Color.RGB = _COL_SUB
            if img_path and img_path.exists():
                try:
                    slide.Shapes.AddPicture(str(img_path.resolve()), LinkToFile=False, SaveWithDocument=True, Left=int(8.8 * 72), Top=int(1.25 * 72), Width=int(4.0 * 72), Height=int(5.5 * 72))
                    print_info(f"    Image inserted: {img_path.name}")
                except Exception as e:
                    print_err(f"    Image failed ({img_path.name}): {e}")
    prs.SaveAs(str(fname.resolve()))
    print_ok(f"PowerPoint saved via COM  ->  {fname.name}")
    return fname, pptApp, prs

def create_docx_via_com(topic: str, llm_text: str, images: list):
    if not WIN32_OK:
        raise RuntimeError("pywin32 is not installed.")
    fname = OUTPUTS / f"{_safe_filename(topic)}_{_timestamp()}.docx"
    pythoncom.CoInitialize()
    print_info("Launching Microsoft Word via COM automation ...")
    wordApp = win32com.client.Dispatch("Word.Application")
    wordApp.Visible = True
    wordApp.WindowState = 1
    doc = wordApp.Documents.Add()
    def _append_para(text: str, style_name: str = "Normal", size: int = 11, bold: bool = False, color_long: int = 0x000000):
        rng = doc.Content
        rng.Collapse(0)
        rng.InsertParagraphAfter()
        rng.Collapse(0)
        rng.Text = text
        rng.Font.Size = size
        rng.Font.Bold = bold
        rng.Font.Color = color_long
        try:
            rng.Style = doc.Styles(style_name)
        except Exception:
            pass
    first_para = doc.Paragraphs(1).Range
    first_para.Text = topic.strip().title()
    try:
        first_para.Style = doc.Styles("Title")
    except Exception:
        pass
    first_para.Font.Size = 28
    first_para.Font.Bold = True
    first_para.Font.Color = _COL_ACCENT
    _append_para(f"Generated by PaperReady  |  Designed by Bhargav  |  Phi-3 Mini (Microsoft)  |  {datetime.now().strftime('%d %b %Y %H:%M')}", size=9, color_long=_rgb_long(0x88, 0x88, 0xAA))
    _append_para("")
    heading_re = re.compile(r"^(?:HEADING\s*[:\-–]|#{1,4})\s*(.+)", re.I)
    img_iter = iter(images)
    first_image = True
    for raw_line in llm_text.splitlines():
        line = raw_line.strip()
        if not line:
            _append_para("")
            continue
        m = heading_re.match(line)
        if m:
            heading_text = m.group(1).strip()
            _append_para(heading_text, style_name="Heading 1", size=14, bold=True, color_long=_COL_ACCENT)
            if first_image:
                img_path = next(img_iter, None)
                if img_path and img_path.exists():
                    try:
                        rng = doc.Content
                        rng.Collapse(0)
                        rng.InsertParagraphAfter()
                        rng.Collapse(0)
                        rng.InlineShapes.AddPicture(FileName=str(img_path.resolve()), LinkToFile=False, SaveWithDocument=True)
                        _append_para("")
                        first_image = False
                        print_info(f"    Image inserted in Word: {img_path.name}")
                    except Exception as e:
                        print_err(f"    Word image failed: {e}")
        else:
            _append_para(line)
    for img_path in img_iter:
        if img_path.exists():
            try:
                rng = doc.Content
                rng.Collapse(0)
                rng.InsertParagraphAfter()
                rng.Collapse(0)
                rng.InlineShapes.AddPicture(FileName=str(img_path.resolve()), LinkToFile=False, SaveWithDocument=True)
            except Exception:
                pass
    doc.SaveAs2(str(fname.resolve()))
    print_ok(f"Word document saved via COM  ->  {fname.name}")
    return fname, wordApp, doc

def create_pptx_fallback(topic: str, llm_text: str, images: list, slide_count: int) -> Path:
    if not PPTX_OK:
        raise RuntimeError("python-pptx not installed.  Run:  pip install python-pptx")
    slides_data = parse_slides(llm_text, topic, slide_count)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    DARK_BG = RGBColor(0x12, 0x12, 0x1E)
    ACCENT = RGBColor(0x6C, 0x63, 0xFF)
    LIGHT_TXT = RGBColor(0xF0, 0xF0, 0xFF)
    SUB_TXT = RGBColor(0xA0, 0x9E, 0xC5)
    TITLE_CLR = RGBColor(0x00, 0xE5, 0xCC)
    BULLET_CLR = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    COVER_TITLE_CLR = RGBColor(0x29, 0xB6, 0xFF)
    blank_layout = prs.slide_layouts[6]
    def hex_bg(slide, color):
        bg = slide.background; fill = bg.fill
        fill.solid(); fill.fore_color.rgb = color
    def add_tb(slide, text, l, t, w, h, fs=18, bold=False, color=None, align=PP_ALIGN.LEFT, bg_color=None):
        if color is None: color = LIGHT_TXT
        txBox = slide.shapes.add_textbox(l, t, w, h)
        tf = txBox.text_frame; tf.word_wrap = True
        if bg_color is not None:
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = bg_color
            txBox.line.fill.background()
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run()
        run.text = text; run.font.size = Pt(fs)
        run.font.bold = bold; run.font.color.rgb = color
        return txBox
    def add_rect(slide, l, t, w, h, color):
        shape = slide.shapes.add_shape(1, l, t, w, h)
        shape.fill.solid(); shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    def add_bullets(slide, bullets, l, t, w, h):
        txBox = slide.shapes.add_textbox(l, t, w, h)
        tf = txBox.text_frame; tf.word_wrap = True
        txBox.fill.solid(); txBox.fill.fore_color.rgb = WHITE
        txBox.line.fill.background()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = f"  \u25b8  {b}"
            run.font.size = Pt(15)
            run.font.color.rgb = BULLET_CLR
        return txBox
    img_iter = iter(images)
    total = len(slides_data)
    for idx, (title, bullets) in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        hex_bg(slide, DARK_BG)
        if idx == 0:
            add_tb(slide, "PaperReady  |  Designed by Bhargav", Inches(0.5), Inches(0.3), Inches(12), Inches(0.5), fs=11, color=ACCENT)
            add_tb(slide, topic.strip().title(), Inches(1), Inches(2.2), Inches(11.3), Inches(2.2), fs=44, bold=True, color=COVER_TITLE_CLR, align=PP_ALIGN.CENTER)
            add_tb(slide, f"{topic.strip().title()}\n{datetime.now().strftime('%B %d, %Y')}", Inches(1), Inches(4.6), Inches(11.3), Inches(1.0), fs=18, color=SUB_TXT, align=PP_ALIGN.CENTER)
            add_rect(slide, Inches(3.5), Inches(5.0), Inches(6.3), Pt(3), ACCENT)
        else:
            add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), WHITE)
            add_rect(slide, Inches(0), Inches(0), Inches(13.33), Pt(5), ACCENT)
            add_tb(slide, title, Inches(0.5), Inches(0.05), Inches(11.5), Inches(1.0), fs=26, bold=True, color=TITLE_CLR)
            add_tb(slide, f"{idx} / {total - 1}", Inches(12.0), Inches(0.2), Inches(1.1), Inches(0.5), fs=11, color=ACCENT, align=PP_ALIGN.RIGHT)
            add_tb(slide, "PaperReady  |  Designed by Bhargav  |  Phi-3 Mini", Inches(0.5), Inches(7.1), Inches(12), Inches(0.35), fs=9, color=SUB_TXT)
            img_path = next(img_iter, None)
            content_width = Inches(8.0) if (img_path and img_path.exists()) else Inches(12.3)
            if bullets:
                add_bullets(slide, bullets, Inches(0.5), Inches(1.3), content_width, Inches(5.4))
            else:
                add_tb(slide, "No content generated for this slide.", Inches(0.5), Inches(1.3), content_width, Inches(5.4), fs=15, color=SUB_TXT)
            if img_path and img_path.exists():
                try:
                    slide.shapes.add_picture(str(img_path), Inches(8.8), Inches(1.25), Inches(4.0), Inches(5.5))
                    print_info(f"Image inserted slide {idx}: {img_path.name}")
                except Exception as e:
                    print_err(f"Image insert failed: {e}")
    fname = OUTPUTS / f"{_safe_filename(topic)}_{_timestamp()}.pptx"
    prs.save(str(fname))
    return fname

def create_docx_fallback(topic: str, llm_text: str, images: list) -> Path:
    if not DOCX_OK:
        raise RuntimeError("python-docx not installed.  Run:  pip install python-docx")
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"; style.font.size = DocPt(11)
    tp = doc.add_heading(topic.strip().title(), level=0)
    if tp.runs:
        tp.runs[0].font.color.rgb = DocRGB(0x6C, 0x63, 0xFF)
        tp.runs[0].font.bold = True
        tp.runs[0].font.size = DocPt(28)
    meta = doc.add_paragraph(f"Generated by PaperReady  |  Designed by Bhargav  |  Phi-3 Mini (Microsoft)  |  {datetime.now().strftime('%d %b %Y %H:%M')}")
    meta.runs[0].font.color.rgb = DocRGB(0x88, 0x88, 0xAA)
    meta.runs[0].font.size = DocPt(9)
    doc.add_paragraph("")
    img_iter = iter(images)
    heading_re = re.compile(r"^(?:HEADING\s*[:\-–]|#{1,4})\s*(.+)", re.I)
    first_image_added = False
    for line in llm_text.splitlines():
        ls = line.strip()
        if not ls:
            doc.add_paragraph(""); continue
        m = heading_re.match(ls)
        if m:
            h = doc.add_heading(m.group(1).strip(), level=1)
            if h.runs: h.runs[0].font.color.rgb = DocRGB(0x6C, 0x63, 0xFF)
            if not first_image_added:
                img_path = next(img_iter, None)
                if img_path and img_path.exists():
                    try:
                        doc.add_picture(str(img_path), width=DocInches(5))
                        doc.add_paragraph(""); first_image_added = True
                        print_info(f"Image in document: {img_path.name}")
                    except Exception as e:
                        print_err(f"Image failed: {e}")
        else:
            p = doc.add_paragraph(ls); p.paragraph_format.space_after = DocPt(6)
    for img_path in img_iter:
        if img_path.exists():
            try:
                doc.add_paragraph("")
                doc.add_picture(str(img_path), width=DocInches(5))
            except Exception:
                pass
    fname = OUTPUTS / f"{_safe_filename(topic)}_{_timestamp()}.docx"
    doc.save(str(fname))
    return fname

def create_txt(topic: str, llm_text: str) -> Path:
    rule = "=" * (W - LEFT)
    header = (
        f"{rule}\n"
        f"  {topic}\n"
        f"  Generated : {datetime.now().strftime('%d %b %Y %H:%M')}\n"
        f"  PaperReady Editor  |  Designed by Bhargav  |  Phi-3 Mini (Microsoft)\n"
        f"{rule}\n\n"
    )
    content = header + llm_text.strip() + "\n"
    fname = OUTPUTS / f"{_safe_filename(topic)}_{_timestamp()}.txt"
    fname.write_text(content, encoding="utf-8")
    return fname
